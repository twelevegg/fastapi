import os
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
import uuid

def load_and_chunk_files(file_paths):
    """Load PDF/PPTX and split into chunk units.

    Returned units keep a minimal citation trail via metadata fields:
      - source: original filename
      - page: PDF page number (from loader metadata) or PPT slide number (1-indexed)
    """

    documents = []
    for path in file_paths:
        if path.lower().endswith('.pdf'):
            loader = PyPDFLoader(path)
            docs = loader.load()
            for d in docs:
                # PyPDFLoader usually provides d.metadata["page"] (0-indexed)
                documents.append({
                    "page_content": d.page_content,
                    "metadata": {
                        "source": path,
                        "page": d.metadata.get("page")
                    }
                })
        elif path.lower().endswith('.pptx'):
            prs = Presentation(path)
            for idx, slide in enumerate(prs.slides, start=1):
                text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                documents.append({
                    "page_content": text,
                    "metadata": {
                        "source": path,
                        "page": idx
                    }
                })
    
    splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
    units = []
    for doc in documents:
        for chunk in splitter.split_text(doc['page_content']):
            if len(chunk.strip()) > 30:
                units.append({
                    "id": str(uuid.uuid4()),
                    "content": chunk,
                    "source": doc['metadata'].get('source'),
                    "page": doc['metadata'].get('page'),
                    "is_learned": False
                })
    return units