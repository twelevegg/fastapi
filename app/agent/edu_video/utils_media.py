import os
import platform
from pathlib import Path
import textwrap
from PIL import Image, ImageDraw, ImageFont
from gtts import gTTS
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
from pptx import Presentation
from pptx.util import Pt


def _pick_font_path():
    """
    OS별로 '실제로 존재하는' 폰트 경로/이름을 찾아 반환합니다.
    - NotoSansKR(프로젝트 번들/시스템 설치) 우선
    - 없으면 OS 기본 한글 폰트로 fallback
    - 최종 실패 시 None
    """
    system = platform.system()  # "Windows", "Darwin", "Linux"

    # 0) (권장) 프로젝트에 번들로 넣은 NotoSansKR 우선 사용
    # 예: 현재 파일 기준 ./assets/fonts/NotoSansKR-Regular.otf
    here = Path(__file__).resolve().parent
    bundled_candidates = [
        here / "assets" / "fonts" / "NotoSansKR-Regular.otf",
        here / "assets" / "fonts" / "NotoSansKR-Regular.ttf",
        here / "assets" / "fonts" / "NotoSansKR-Medium.otf",
        here / "assets" / "fonts" / "NotoSansKR-Bold.otf",
    ]
    for p in bundled_candidates:
        if p.exists():
            return str(p)

    # 1) OS별 시스템 폰트 후보
    if system == "Windows":
        candidates = [
            # NotoSansKR (설치돼 있다면)
            r"C:\Windows\Fonts\NotoSansKR-Regular.otf",
            r"C:\Windows\Fonts\NotoSansKR-Regular.ttf",
            # 맑은 고딕
            "malgun.ttf",
            r"C:\Windows\Fonts\malgun.ttf",
            r"C:\Windows\Fonts\malgunbd.ttf",
        ]
    elif system == "Darwin":  # macOS
        candidates = [
            # NotoSansKR (설치돼 있다면)
            "/Library/Fonts/NotoSansKR-Regular.otf",
            "/Library/Fonts/NotoSansKR-Regular.ttf",
            "/System/Library/Fonts/Supplemental/NotoSansKR-Regular.otf",
            "/System/Library/Fonts/Supplemental/NotoSansKR-Regular.ttf",
            # macOS 기본 한글 폰트들
            "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
            "/Library/Fonts/AppleGothic.ttf",
            "/System/Library/Fonts/AppleSDGothicNeo.ttc",
        ]
    else:  # Linux
        candidates = [
            # Noto Sans KR / CJK (배포판마다 위치가 다를 수 있음)
            "/usr/share/fonts/opentype/noto/NotoSansKR-Regular.otf",
            "/usr/share/fonts/truetype/noto/NotoSansKR-Regular.otf",
            "/usr/share/fonts/truetype/noto/NotoSansKR-Regular.ttf",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            # 대체 한글 폰트
            "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        ]

    for p in candidates:
        try:
            if Path(p).exists():
                return p
        except Exception:
            # "malgun.ttf"처럼 파일명만 있는 경우 Path.exists가 의미 없을 수 있으므로
            # 아래 truetype 테스트에서 걸러지게 둠
            pass

    # 2) Windows에서 "malgun.ttf" 처럼 이름만 주는 경우를 위해 마지막으로 반환 시도
    # (실제 존재 여부는 truetype에서 판단)
    if system == "Windows":
        return "malgun.ttf"

    return None


def generate_ppt_image(text, title, brand, filename):
    if isinstance(text, list):
        text = "\n".join(text)
    text, title, brand = str(text), str(title), str(brand)

    width, height = 1280, 720
    img = Image.new("RGB", (width, height), color=(248, 249, 250))
    d = ImageDraw.Draw(img)

    # 테마 색상 (기업마다 다를 수 있으나 신뢰감을 주는 Deep Blue/Gray 계열로 범용 설정)
    primary_color = (33, 37, 41)
    point_color = (0, 102, 204)  # 포인트 바 색상

    # 폰트 로딩 (OS별 자동 선택)
    try:
        font_path = _pick_font_path()
        if not font_path:
            raise FileNotFoundError("No usable font found")

        # (중요) truetype 로딩 시도 (malgun.ttf처럼 이름만 줘도 Windows에선 로드됨)
        title_font = ImageFont.truetype(font_path, 45)
        body_font = ImageFont.truetype(font_path, 24)
        brand_font = ImageFont.truetype(font_path, 30)
    except Exception:
        title_font = body_font = brand_font = ImageFont.load_default()

    # 1. 상단 제목 바
    d.rectangle([0, 0, 1280, 100], fill=primary_color)
    d.text((40, 30), title, fill=(255, 255, 255), font=title_font)

    # 2. 브랜드 로고 (왼쪽 하단) - 현재는 주석 처리
    # d.text((40, 645), brand, fill=point_color, font=brand_font)
    # d.line([(40, 685), (150, 685)], fill=point_color, width=3)

    # 3. 본문 불렛포인트 출력
    y_start = 130
    lines = text.split("\n")
    for line in lines:
        clean_line = line.strip()
        if not clean_line:
            continue
        if not clean_line.startswith("•"):
            clean_line = "• " + clean_line

        wrapped = textwrap.wrap(clean_line, width=65)
        box_h = (len(wrapped) * 35) + 15

        if y_start + box_h > 630:
            break

        d.rounded_rectangle(
            [60, y_start, 1220, y_start + box_h],
            radius=5,
            fill=(255, 255, 255),
            outline=(230, 230, 230),
        )
        d.rectangle([60, y_start + 5, 65, y_start + box_h - 5], fill=point_color)

        curr_y = y_start + 8
        for w_line in wrapped:
            d.text((85, curr_y), w_line, fill=(50, 50, 50), font=body_font)
            curr_y += 35
        y_start += box_h + 10

    img.save(filename)


def create_video_segment(script_segments, output_filename="output.mp4"):
    clips = []
    for i, seg in enumerate(script_segments):
        img_path = f"temp_{i}.png"
        generate_ppt_image(seg["summary"], seg["title"], seg.get("brand", "Education"), img_path)

        audio_path = f"temp_{i}.mp3"
        tts = gTTS(text=seg["text"], lang="ko")
        tts.save(audio_path)

        audio_clip = AudioFileClip(audio_path)
        video_clip = ImageClip(img_path).set_duration(audio_clip.duration + 0.5)
        video_clip = video_clip.set_audio(audio_clip)
        clips.append(video_clip)

    final_video = concatenate_videoclips(clips, method="compose")
    final_video.write_videofile(output_filename, fps=24, codec="libx264")

    for i in range(len(script_segments)):
        try:
            os.remove(f"temp_{i}.png")
            os.remove(f"temp_{i}.mp3")
        except:
            pass


def create_pptx(script_segments, output_filename="output.pptx"):
    """Create a simple PPTX from generated script segments.

    - Slide title: seg['title']
    - Bullets: seg['summary'] split by lines
    - Speaker notes: seg['text']
    """
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(title_slide_layout)
    s0.shapes.title.text = "교육 세션"
    subtitle = s0.placeholders[1]
    brand = (script_segments[0].get("brand") if script_segments else "")
    subtitle.text = brand if brand else ""

    # Content slides
    content_layout = prs.slide_layouts[1]
    for seg in script_segments:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = str(seg.get("title", ""))

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        summary = seg.get("summary", "")
        lines = [ln.strip() for ln in str(summary).split("\n") if ln.strip()]
        if not lines:
            lines = ["(요약 내용 없음)"]

        for j, ln in enumerate(lines):
            ln = ln.lstrip("• ").strip()
            p = body.paragraphs[0] if j == 0 else body.add_paragraph()
            p.text = ln
            p.level = 0
            p.font.size = Pt(20)

        # Speaker notes
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.text = str(seg.get("text", ""))

    prs.save(output_filename)
    return output_filename
